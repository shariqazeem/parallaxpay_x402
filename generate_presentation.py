#!/usr/bin/env python3
"""
ParallaxPay Presentation Generator
Generates a professional PowerPoint presentation for x402 Solana Hackathon submission
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate the complete ParallaxPay presentation"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 aspect ratio
    prs.slide_height = Inches(9)

    # Define color scheme (Solana-inspired)
    PURPLE = RGBColor(138, 43, 226)  # Purple
    BLUE = RGBColor(20, 241, 149)    # Solana green-blue
    DARK_BG = RGBColor(18, 18, 18)   # Dark background
    LIGHT_TEXT = RGBColor(255, 255, 255)  # White text
    ACCENT_GREEN = RGBColor(0, 255, 163)  # Green accent

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dark background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ParallaxPay"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(96)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Crypto Oracle with x402 Micropayments"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = LIGHT_TEXT

    # Hackathon info
    info_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(0.8))
    info_frame = info_box.text_frame
    info_frame.text = "x402 Solana Hackathon - Parallax Eco Track"
    info_para = info_frame.paragraphs[0]
    info_para.alignment = PP_ALIGN.CENTER
    info_para.font.size = Pt(24)
    info_para.font.color.rgb = ACCENT_GREEN

    # Author
    author_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(0.5))
    author_frame = author_box.text_frame
    author_frame.text = "Shariq Azeem"
    author_para = author_frame.paragraphs[0]
    author_para.alignment = PP_ALIGN.CENTER
    author_para.font.size = Pt(28)
    author_para.font.color.rgb = LIGHT_TEXT

    # Screenshot placeholder
    screenshot_box = slide1.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(0.8))
    screenshot_frame = screenshot_box.text_frame
    screenshot_frame.text = "[Add screenshot: Landing page hero from https://parallaxpay.online]"
    screenshot_para = screenshot_frame.paragraphs[0]
    screenshot_para.alignment = PP_ALIGN.CENTER
    screenshot_para.font.size = Pt(18)
    screenshot_para.font.italic = True
    screenshot_para.font.color.rgb = RGBColor(150, 150, 150)

    # Speaker notes
    notes_slide = slide1.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Hi, I'm Shariq Azeem, and I built ParallaxPay for the x402 Solana Hackathon. This is a decentralized AI oracle that predicts cryptocurrency prices using multi-provider consensus and x402 micropayments."

    # SLIDE 2: The Problem
    slide2 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="The Problem with Traditional Oracles",
        bullets=[
            "💸 Traditional APIs: $99/month for basic coverage",
            "🔒 Limited to 4-5 major cryptocurrencies",
            "🏢 Centralized providers (single point of failure)",
            "⚠️ No transparency in prediction accuracy"
        ],
        notes="Today's crypto oracles have major problems. They charge $99 per month, only cover Bitcoin and Ethereum, rely on centralized providers, and give you zero transparency. If you want to predict smaller altcoins or DeFi tokens? You're out of luck."
    )

    # SLIDE 3: The Solution
    slide3 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="ParallaxPay Solution",
        bullets=[
            "🪙 152+ Cryptocurrencies - From Bitcoin to meme coins",
            "🤖 Multi-Provider AI Consensus - Powered by Gradient Parallax",
            "💰 x402 Micropayments - Pay $0.0008 per query (92% savings)",
            "🔄 Autonomous Operation - 24/7 automated predictions"
        ],
        notes="ParallaxPay solves all of this. You get 152 different cryptocurrencies, multi-provider AI consensus through Gradient Parallax, x402 micropayments at less than a penny per query - that's 92% cost savings - and it runs autonomously 24/7.",
        screenshot_note="[Add screenshot: Oracle page with coin selector dropdown showing many coins]"
    )

    # SLIDE 4: Live Demo - Oracle Agent
    slide4 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="Market Oracle in Action",
        bullets=[
            "✓ Select from 152+ cryptocurrencies",
            "✓ Multi-provider AI consensus voting",
            "✓ Cost: $0.0008 per prediction",
            "✓ 85%+ confidence scores",
            "✓ Real-time market analysis"
        ],
        notes="Let me show you how it works. Here I'm selecting BNB, one of 152 supported coins. Click 'Make Prediction' and watch - three AI providers from the Parallax network analyze the market and vote on the prediction. This is multi-provider consensus in action. The result? A prediction with 85% confidence that cost $0.0008. Traditional APIs would charge you a monthly subscription for just a fraction of this coverage.",
        screenshot_note="[Add 3 screenshots: 1) Oracle page with BNB selected, 2) Prediction loading, 3) Result with cost]"
    )

    # SLIDE 5: Decentralized Marketplace
    slide5 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="Open AI Marketplace",
        bullets=[
            "🌐 Real Parallax nodes competing for requests",
            "📊 Performance metrics & reputation scores",
            "🚀 Anyone can become a provider",
            "💎 Best providers get rewarded",
            "⚡ Permissionless participation"
        ],
        notes="What makes this truly decentralized? The marketplace. These are real Parallax nodes running on the network, competing for requests based on performance. Anyone can become a provider, and the best ones earn rewards. It's a true permissionless AI marketplace.",
        screenshot_note="[Add screenshot: Marketplace page showing provider cards with stats]"
    )

    # SLIDE 6: Analytics & Transparency
    slide6 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="Complete Transparency",
        bullets=[
            "📈 Real-time cost comparison charts",
            "💰 92% savings vs traditional APIs",
            "🎯 Prediction accuracy tracking",
            "🔗 Full x402 transaction history",
            "✅ Every interaction on-chain"
        ],
        notes="Every interaction is tracked on-chain. The analytics dashboard shows your cost savings - 92% compared to traditional APIs - prediction accuracy over time, and complete transaction history. Full transparency powered by Solana and x402.",
        screenshot_note="[Add screenshots: Analytics dashboard + Transaction history page]"
    )

    # SLIDE 7: Technical Architecture
    slide7 = add_content_slide(prs, DARK_BG, PURPLE, LIGHT_TEXT, ACCENT_GREEN,
        title="Built on Cutting-Edge Tech",
        bullets=[
            "⚛️ Frontend: Next.js + TypeScript + TailwindCSS",
            "⛓️ Blockchain: Solana (x402 protocol)",
            "🤖 AI Layer: Gradient Parallax distributed inference",
            "💾 Data: Supabase PostgreSQL",
            "💸 Payments: x402 micropayments (sub-cent)",
            "🐳 Deployment: Docker + Oracle Cloud"
        ],
        notes="Here's the tech stack: Next.js frontend, Solana with x402 for micropayments, Gradient Parallax for distributed AI inference, and Supabase for data persistence. The entire platform is containerized and deployed on Oracle Cloud. Users can even deploy their own custom agents with one click from the agent builder.",
        screenshot_note="[Add screenshot: Agent builder page showing deployment options]"
    )

    # SLIDE 8: Why ParallaxPay Wins
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide8.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Why ParallaxPay?"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = PURPLE

    # Competitive advantages - Large text
    advantages = [
        "✅ 152+ cryptocurrencies vs competitors' 4-5",
        "✅ Real Parallax integration (not simulated)",
        "✅ Real x402 micropayments (on-chain verification)",
        "✅ 92% cost savings with transparent pricing",
        "✅ Autonomous mode for 24/7 operation",
        "✅ Production-ready with live deployment"
    ]

    content_box = slide8.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for advantage in advantages:
        p = content_frame.add_paragraph()
        p.text = advantage
        p.font.size = Pt(32)
        p.font.color.rgb = LIGHT_TEXT
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

    # URL
    url_box = slide8.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(0.8))
    url_frame = url_box.text_frame
    url_frame.text = "🌐 https://parallaxpay.online"
    url_para = url_frame.paragraphs[0]
    url_para.alignment = PP_ALIGN.CENTER
    url_para.font.size = Pt(36)
    url_para.font.bold = True
    url_para.font.color.rgb = ACCENT_GREEN

    # Closing statement
    closing_box = slide8.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
    closing_frame = closing_box.text_frame
    closing_frame.text = "Infrastructure for the Future of AI Agents on Solana"
    closing_para = closing_frame.paragraphs[0]
    closing_para.alignment = PP_ALIGN.CENTER
    closing_para.font.size = Pt(28)
    closing_para.font.color.rgb = LIGHT_TEXT
    closing_para.font.italic = True

    # Speaker notes
    notes_slide = slide8.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Why does ParallaxPay win? We support 152 cryptocurrencies versus competitors' 4 or 5. This is real Parallax integration with real x402 micropayments - fully on-chain and verifiable. Users save 92%, run agents autonomously, and it's production-ready right now at parallaxpay.online. This isn't a proof of concept - it's infrastructure for the future of AI agents on Solana."

    # Save presentation
    filename = "ParallaxPay_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"\n📸 Next steps:")
    print(f"1. Open {filename} in PowerPoint")
    print(f"2. Add screenshots from https://parallaxpay.online (look for placeholder notes)")
    print(f"3. Review speaker notes (visible in Notes pane)")
    print(f"4. Practice once, then record (Slide Show → Record Slide Show)")
    print(f"5. Export as video (File → Export → Create Video)")
    print(f"\n🚀 You're ready to submit!")

    return filename


def add_content_slide(prs, bg_color, title_color, text_color, accent_color, title, bullets, notes, screenshot_note=None):
    """Add a content slide with title, bullets, and notes"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color

    # Content bullets
    content_top = Inches(1.8)
    content_height = Inches(5) if not screenshot_note else Inches(4.5)

    content_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(14), content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(32)
        p.font.color.rgb = text_color
        p.space_before = Pt(18)
        p.space_after = Pt(18)
        p.level = 0

    # Screenshot note if provided
    if screenshot_note:
        screenshot_box = slide.shapes.add_textbox(Inches(1), Inches(7.8), Inches(14), Inches(0.8))
        screenshot_frame = screenshot_box.text_frame
        screenshot_frame.text = screenshot_note
        screenshot_para = screenshot_frame.paragraphs[0]
        screenshot_para.alignment = PP_ALIGN.CENTER
        screenshot_para.font.size = Pt(16)
        screenshot_para.font.italic = True
        screenshot_para.font.color.rgb = RGBColor(150, 150, 150)

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes

    return slide


if __name__ == "__main__":
    print("🎯 Generating ParallaxPay Hackathon Presentation...")
    print("=" * 60)
    create_presentation()
